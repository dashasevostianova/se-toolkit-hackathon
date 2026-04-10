"""Generate presentation PowerPoint for DailyEnglish — Updated Version 2."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import qrcode
import io
from PIL import Image

# ===== COLORS =====
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT_LIGHT = RGBColor(0xA2, 0x9B, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB8)
MUTED = RGBColor(0x6B, 0x6B, 0x80)
SUCCESS = RGBColor(0x00, 0xCE, 0xC9)
WARNING = RGBColor(0xFD, 0xCB, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet(text_frame, text, font_size=16, color=GRAY, bullet="\u2022", space_before=Pt(8)):
    p = text_frame.add_paragraph()
    p.text = f"{bullet}  {text}"
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_before = space_before
    return p


def make_qr(url):
    qr = qrcode.QRCode(version=1, box_size=10, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="white", back_color="black")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)

# Decorative shapes
shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-0.5), Inches(4), Inches(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
shape.shadow.inherit = False

shape2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SUCCESS
shape2.line.fill.background()
shape2.shadow.inherit = False

# Logo + Title
add_text_box(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
    "\U0001f4da  DailyEnglish", font_size=52, color=WHITE, bold=True)

add_text_box(slide1, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
    "Learn one English word at a time — automatically", font_size=28, color=ACCENT_LIGHT, bold=False)

# Divider
add_shape(slide1, Inches(1.5), Inches(3.8), Inches(3), Inches(0.03), ACCENT)

# Author info
add_text_box(slide1, Inches(1.5), Inches(4.2), Inches(6), Inches(0.4),
    "Dasha Sevostianova", font_size=20, color=WHITE, bold=True)
add_text_box(slide1, Inches(1.5), Inches(4.7), Inches(6), Inches(0.4),
    "Group: CSE-02", font_size=18, color=GRAY)
add_text_box(slide1, Inches(1.5), Inches(5.1), Inches(6), Inches(0.4),
    "Email: dasha.sevostianova@university.ru", font_size=18, color=GRAY)


# ============================================================
# SLIDE 2 — CONTEXT
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
    "Context", font_size=36, color=ACCENT_LIGHT, bold=True)
add_shape(slide2, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT)

col_w = Inches(3.5)
col_h = Inches(4.5)
gap = Inches(0.5)
start_x = Inches(0.8)
card_y = Inches(1.8)

# Card 1 — End User
add_shape(slide2, start_x, card_y, col_w, col_h, BG_CARD, RGBColor(0x2A, 0x2A, 0x40), 1)
txBox = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(3.2), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "\U0001f465  End-User"
p.font.size = Pt(22)
p.font.color.rgb = WHITE
p.font.bold = True
add_bullet(tf, "Anyone from Russia learning English", 17, GRAY, space_before=Pt(12))
add_bullet(tf, "Students, professionals, casual learners", 17, GRAY)
add_bullet(tf, "Prefers lightweight, no-setup tools", 17, GRAY)
add_bullet(tf, "Wants regular, frictionless practice", 17, GRAY)

# Card 2 — Problem
add_shape(slide2, start_x + col_w + gap, card_y, col_w, col_h, BG_CARD, RGBColor(0x2A, 0x2A, 0x40), 1)
txBox2 = slide2.shapes.add_textbox(Inches(4.8), Inches(2.1), Inches(3.2), Inches(4))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "\u26a0\ufe0f  Problem"
p2.font.size = Pt(22)
p2.font.color.rgb = WHITE
p2.font.bold = True
add_bullet(tf2, "Apps require manual opening", 17, GRAY, space_before=Pt(12))
add_bullet(tf2, "Overloaded with features", 17, GRAY)
add_bullet(tf2, "No automatic reminders", 17, GRAY)
add_bullet(tf2, "Motivation drops quickly", 17, GRAY)
add_bullet(tf2, "\u2192 Irregular, inconsistent practice", 17, RGBColor(0xFF, 0x6B, 0x6B))

# Card 3 — Solution
add_shape(slide2, start_x + 2 * (col_w + gap), card_y, col_w, col_h, BG_CARD, RGBColor(0x2A, 0x2A, 0x40), 1)
txBox3 = slide2.shapes.add_textbox(Inches(8.8), Inches(2.1), Inches(3.5), Inches(4))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "\U0001f4a1  Solution"
p3.font.size = Pt(22)
p3.font.color.rgb = WHITE
p3.font.bold = True
p_idea = tf3.add_paragraph()
p_idea.text = "A web app that auto-serves one new English word every 3 minutes — with translation, example, and self-test mode."
p_idea.font.size = Pt(17)
p_idea.font.color.rgb = SUCCESS
p_idea.font.name = "Calibri"
p_idea.space_before = Pt(12)
add_bullet(tf3, "No login, no setup needed", 17, GRAY)
add_bullet(tf3, "User sets daily limit (10/20/all)", 17, GRAY)


# ============================================================
# SLIDE 3 — IMPLEMENTATION
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
    "Implementation", font_size=36, color=ACCENT_LIGHT, bold=True)
add_shape(slide3, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT)

# Left — Final Version
add_shape(slide3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), BG_CARD, RGBColor(0x2A, 0x2A, 0x40), 1)
txV1 = slide3.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5), Inches(4.8))
tfV1 = txV1.text_frame
tfV1.word_wrap = True
p = tfV1.paragraphs[0]
p.text = "\U0001f310  Version 2 — Web App (Final)"
p.font.size = Pt(24)
p.font.color.rgb = SUCCESS
p.font.bold = True
add_bullet(tfV1, "Flask backend + REST API", 17, GRAY, space_before=Pt(12))
add_bullet(tfV1, "Dark animated UI (CSS + JS)", 17, GRAY)
add_bullet(tfV1, "Auto-refresh every 3 minutes", 17, GRAY)
add_bullet(tfV1, "Daily word limit: 10 / 20 / 90", 17, GRAY)
add_bullet(tfV1, "Translation hidden \u2192 click to reveal", 17, GRAY)
add_bullet(tfV1, "Random order, reshuffle after pool", 17, GRAY)
add_bullet(tfV1, "Progress bar + round counter", 17, GRAY)

# Divider
add_shape(slide3, Inches(6.6), Inches(2), Inches(0.03), Inches(4.5), RGBColor(0x2A, 0x2A, 0x40))

# Right — Tech + Feedback
add_shape(slide3, Inches(7), Inches(1.6), Inches(5.5), Inches(5.2), BG_CARD, RGBColor(0x2A, 0x2A, 0x40), 1)
txV2 = slide3.shapes.add_textbox(Inches(7.3), Inches(1.9), Inches(5), Inches(4.8))
tfV2 = txV2.text_frame
tfV2.word_wrap = True
p = tfV2.paragraphs[0]
p.text = "\U0001f6e0\ufe0f  Tech Stack + Feedback"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.font.bold = True
add_bullet(tfV2, "Python 3.12 + Flask 3.0", 17, GRAY, space_before=Pt(12))
add_bullet(tfV2, "Vanilla HTML/CSS/JS (no frameworks)", 17, GRAY)
add_bullet(tfV2, "Deployed via nohup / systemd", 17, GRAY)
add_bullet(tfV2, "Optional Docker Compose support", 17, GRAY)
add_bullet(tfV2, "90-word curated dictionary", 17, GRAY)
add_bullet(tfV2, "", font_size=8, space_before=Pt(4))
add_bullet(tfV2, "\u2705  TA Feedback Addressed:", font_size=15, color=SUCCESS, space_before=Pt(6))
add_bullet(tfV2, "Added daily limit selector", 17, SUCCESS)
add_bullet(tfV2, "Words shuffled, not in order", 17, SUCCESS)
add_bullet(tfV2, "Translation hidden for self-test", 17, SUCCESS)


# ============================================================
# SLIDE 4 — DEMO
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
    "Demo", font_size=36, color=ACCENT_LIGHT, bold=True)
add_shape(slide4, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT)

# Video placeholder
add_shape(slide4, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5), BG_CARD, RGBColor(0x2A, 0x2A, 0x40), 2)
add_text_box(slide4, Inches(2), Inches(3), Inches(9.3), Inches(0.8),
    "\U0001f3ac  Pre-recorded Video Demonstration of Version 2",
    font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(2), Inches(4), Inches(9.3), Inches(0.6),
    "Video with voice-over (\u2264 2 minutes)",
    font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide4, Inches(2), Inches(4.7), Inches(9.3), Inches(0.5),
    "[ Embed video file here ]",
    font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

# Key features
add_text_box(slide4, Inches(0.8), Inches(6.6), Inches(12), Inches(0.5),
    "Key Features Shown: Word card with emoji  \u2022  Hidden translation (click to reveal)  \u2022  Daily limit settings  \u2022  Auto-refresh timer  \u2022  Progress tracking  \u2022  Reshuffle notification",
    font_size=14, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — LINKS
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
    "Links", font_size=36, color=ACCENT_LIGHT, bold=True)
add_shape(slide5, Inches(0.8), Inches(1.2), Inches(2), Inches(0.04), ACCENT)

# GitHub QR + Link
gh_url = "https://github.com/dashasevostianova/se-toolkit-hackathon"
qr_buf = make_qr(gh_url)
slide5.shapes.add_picture(qr_buf, Inches(1.5), Inches(2), Inches(2.5), Inches(2.5))

add_text_box(slide5, Inches(4.5), Inches(2.2), Inches(7), Inches(0.5),
    "GitHub Repository", font_size=24, color=WHITE, bold=True)
add_text_box(slide5, Inches(4.5), Inches(2.9), Inches(7), Inches(0.5),
    gh_url, font_size=16, color=ACCENT_LIGHT)

# Divider
add_shape(slide5, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.02), RGBColor(0x2A, 0x2A, 0x40))

# Deployed App QR + Link
deploy_url = "http://10.93.25.94:5000"
qr_buf2 = make_qr(deploy_url)
slide5.shapes.add_picture(qr_buf2, Inches(1.5), Inches(5.2), Inches(2.5), Inches(2.5))

add_text_box(slide5, Inches(4.5), Inches(5.4), Inches(7), Inches(0.5),
    "Deployed Product (Version 2)", font_size=24, color=WHITE, bold=True)
add_text_box(slide5, Inches(4.5), Inches(6.1), Inches(7), Inches(0.5),
    deploy_url, font_size=16, color=SUCCESS)

# Thank you
add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(12), Inches(0.5),
    "Thank you! \U0001f393", font_size=20, color=MUTED, alignment=PP_ALIGN.RIGHT)


# ===== SAVE =====
output_path = r"c:\Users\sevos\OneDrive\Desktop\UI PROGRAMS\semester2\project\DailyEnglish_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
